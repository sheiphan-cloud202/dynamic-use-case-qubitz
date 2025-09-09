#!/usr/bin/env python3
"""
Create PowerPoint template files for AWS deployment
Run this locally to generate template files, then upload to S3
"""
import os
import boto3
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

def create_all_templates():
    """Create all template files for AWS deployment"""
    
    if not PPT_AVAILABLE:
        print("❌ python-pptx not available. Install with: pip install python-pptx")
        return
    
    print("🎨 Creating PowerPoint Templates for AWS Deployment")
    print("=" * 55)
    
    # Template configurations
    templates = {
        "first_deck": {
            "title": "Strategic Partnership Opportunity",
            "subtitle": "[Company Name] Executive Overview",
            "colors": {
                "primary": RGBColor(20, 33, 61),
                "secondary": RGBColor(52, 73, 94),
                "accent": RGBColor(230, 126, 34)
            }
        },
        "marketing": {
            "title": "Transform Your Business Today",
            "subtitle": "Unlock Growth and Innovation",
            "colors": {
                "primary": RGBColor(225, 45, 139),
                "secondary": RGBColor(74, 144, 226),
                "accent": RGBColor(255, 193, 7)
            }
        },
        "use_case": {
            "title": "Use Case Implementation Strategy",
            "subtitle": "Transformation Scenarios",
            "colors": {
                "primary": RGBColor(99, 102, 241),
                "secondary": RGBColor(139, 69, 19),
                "accent": RGBColor(245, 158, 11)
            }
        },
        "technical": {
            "title": "Technical Architecture Overview",
            "subtitle": "System Design and Implementation",
            "colors": {
                "primary": RGBColor(30, 41, 59),
                "secondary": RGBColor(71, 85, 105),
                "accent": RGBColor(14, 165, 233)
            }
        },
        "strategy": {
            "title": "Strategic Transformation Roadmap",
            "subtitle": "3-Year Strategic Plan",
            "colors": {
                "primary": RGBColor(79, 70, 229),
                "secondary": RGBColor(107, 114, 128),
                "accent": RGBColor(16, 185, 129)
            }
        }
    }
    
    created_files = []
    
    for template_name, config in templates.items():
        try:
            print(f"\n📋 Creating {template_name} template...")
            
            # Create presentation
            ppt = Presentation()
            
            # Add title slide
            slide = ppt.slides.add_slide(ppt.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = config["title"]
            subtitle.text = config["subtitle"]
            
            # Apply styling
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.color.rgb = config["colors"]["primary"]
            title.text_frame.paragraphs[0].font.bold = True
            
            subtitle.text_frame.paragraphs[0].font.size = Pt(26)
            subtitle.text_frame.paragraphs[0].font.color.rgb = config["colors"]["secondary"]
            
            # Add sample content slide
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"{template_name.title()} Overview"
            content.text = f"• {template_name.title()} presentation structure\n• Professional styling and layout\n• Ready for content generation\n• Optimized for business use"
            
            # Style content slide
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = config["colors"]["primary"]
            title.text_frame.paragraphs[0].font.bold = True
            
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(55, 65, 81)
                paragraph.space_before = Pt(6)
            
            # Save template
            filename = f"{template_name}_template.pptx"
            ppt.save(filename)
            
            file_size = os.path.getsize(filename)
            print(f"✅ Created: {filename} ({file_size:,} bytes)")
            created_files.append(filename)
            
        except Exception as e:
            print(f"❌ Error creating {template_name}: {e}")
    
    print(f"\n🎉 Template Creation Complete!")
    print(f"Created {len(created_files)} template files:")
    for filename in created_files:
        print(f"  📎 {filename}")
    
    return created_files

def upload_templates_to_s3(bucket_name: str = None):
    """Upload templates to S3 for Lambda access"""
    
    if not bucket_name:
        print("❌ S3 bucket name required for upload")
        return
    
    print(f"\n☁️ Uploading templates to S3 bucket: {bucket_name}")
    
    s3_client = boto3.client('s3')
    
    template_files = [f for f in os.listdir('.') if f.endswith('_template.pptx')]
    
    for filename in template_files:
        try:
            s3_key = f"templates/{filename}"
            s3_client.upload_file(filename, bucket_name, s3_key)
            print(f"✅ Uploaded: {filename} → s3://{bucket_name}/{s3_key}")
        except Exception as e:
            print(f"❌ Failed to upload {filename}: {e}")

if __name__ == "__main__":
    # Create templates
    created_files = create_all_templates()
    
    if created_files:
        print(f"\n📤 Upload to S3? (optional)")
        bucket = input("Enter S3 bucket name (or press Enter to skip): ").strip()
        
        if bucket:
            upload_templates_to_s3(bucket)
        else:
            print("Templates created locally. Upload manually or use in development.")