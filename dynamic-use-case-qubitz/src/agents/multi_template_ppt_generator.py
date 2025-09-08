#!/usr/bin/env python3
"""
Complete Multi-Template PowerPoint Generation Agent
Integrates with the existing orchestrator to provide 5 different presentation types
"""
import os
import boto3
import json
import logging
from datetime import datetime
from typing import Dict, Any, List, Optional
from abc import ABC, abstractmethod

from strands import Agent
from strands_tools import retrieve, http_request
from src.core.bedrock_manager import EnhancedModelManager
from src.core.models import CompanyProfile, UseCaseStructured
from src.services.aws_clients import s3_client, S3_BUCKET, LAMBDA_TMP_DIR
from src.utils.status_tracker import StatusTracker, StatusCheckpoints

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
    logger.info("✅ python-pptx available for PowerPoint generation")
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("⚠️ python-pptx not available")

class PPTTemplate(ABC):
    """Abstract base class for PowerPoint templates"""
    
    def __init__(self, name: str, colors: Dict[str, RGBColor]):
        self.name = name
        self.colors = colors
    
    @abstractmethod
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        """Generate Bedrock prompt specific to this template type"""
        pass
    
    @abstractmethod
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create slides specific to this template type"""
        pass

class FirstDeckTemplate(PPTTemplate):
    """First Deck Call Template - High-level executive overview"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(20, 33, 61),     # Deep navy
            "secondary": RGBColor(52, 73, 94),   # Slate blue
            "accent": RGBColor(230, 126, 34),    # Orange
            "text": RGBColor(44, 62, 80),        # Dark blue-gray
            "light": RGBColor(236, 240, 241)     # Light gray
        }
        super().__init__("First Deck", colors)
    
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        
        # Get top 3 use cases for executive overview
        top_use_cases = use_cases[:3] if use_cases else []
        
        return f"""
Create a FIRST DECK CALL presentation for {company_profile.name}. This is an initial executive presentation.

COMPANY CONTEXT:
- Industry: {company_profile.industry}
- Business Model: {company_profile.business_model}
- Company Size: {company_profile.company_size}

TOP STRATEGIC OPPORTUNITIES:
{self._format_use_cases_for_executive(top_use_cases)}

RESEARCH INSIGHTS:
{research_data.get('research_findings', '')[:1000]}

Create a high-level executive presentation focused on strategic partnership opportunities.

OUTPUT JSON:
{{
    "presentation_title": "Strategic Partnership Opportunity - {company_profile.name}",
    "slides": [
        {{
            "type": "title",
            "title": "Strategic Partnership Opportunity",
            "subtitle": "{company_profile.name} Business Transformation Overview"
        }},
        {{
            "type": "company_overview",
            "title": "Company Overview",
            "key_points": [
                "Business profile in {company_profile.industry}",
                "Current market position and scale",
                "Key strategic objectives and priorities"
            ]
        }},
        {{
            "type": "opportunities",
            "title": "Strategic Opportunities Identified", 
            "opportunities": [
                "Primary transformation opportunity with significant impact",
                "Efficiency optimization potential for competitive advantage",
                "Innovation enablement areas for market leadership"
            ]
        }},
        {{
            "type": "value_proposition",
            "title": "Potential Business Value",
            "value_items": [
                "Cost optimization: 15-25% efficiency potential",
                "Revenue enhancement: Growth acceleration opportunities",
                "Strategic positioning: Competitive advantage development"
            ]
        }},
        {{
            "type": "next_steps",
            "title": "Proposed Next Steps",
            "steps": [
                "Detailed assessment and discovery phase",
                "Proof of concept development and validation", 
                "Implementation roadmap and partnership framework"
            ]
        }}
    ]
}}

Focus on high-level business impact and executive-appropriate messaging.
"""
    
    def _format_use_cases_for_executive(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for executive summary"""
        if not use_cases:
            return "Standard business transformation opportunities"
        
        formatted = []
        for i, uc in enumerate(use_cases[:3], 1):
            formatted.append(f"{i}. {uc.title} - {uc.business_value[:100]}...")
        
        return "\n".join(formatted)
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create first deck slides"""
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_title_slide(ppt, slide_data)
            elif slide_type == 'company_overview':
                self._create_overview_slide(ppt, slide_data)
            elif slide_type == 'opportunities':
                self._create_opportunities_slide(ppt, slide_data)
            elif slide_type == 'value_proposition':
                self._create_value_slide(ppt, slide_data)
            elif slide_type == 'next_steps':
                self._create_next_steps_slide(ppt, slide_data)
    
    def _create_title_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Partnership')
        subtitle.text = data.get('subtitle', 'Executive Overview')
        
        # Executive styling
        title.text_frame.paragraphs[0].font.size = Pt(48)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_overview_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Company Overview')
        
        key_points = data.get('key_points', [])
        if key_points:
            content.text = key_points[0]
            for point in key_points[1:]:
                p = content.text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _create_opportunities_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Opportunities')
        
        opportunities = data.get('opportunities', [])
        if opportunities:
            content.text = opportunities[0]
            for opp in opportunities[1:]:
                p = content.text_frame.add_paragraph()
                p.text = opp
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _create_value_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Business Value')
        
        value_items = data.get('value_items', [])
        if value_items:
            content.text = value_items[0]
            for item in value_items[1:]:
                p = content.text_frame.add_paragraph()
                p.text = item
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _create_next_steps_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Next Steps')
        
        steps = data.get('steps', [])
        if steps:
            content.text = steps[0]
            for step in steps[1:]:
                p = content.text_frame.add_paragraph()
                p.text = step
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _style_content_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(8)

class MarketingTemplate(PPTTemplate):
    """Marketing Template - Persuasive, benefit-focused"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(225, 45, 139),    # Vibrant pink
            "secondary": RGBColor(74, 144, 226),  # Bright blue
            "accent": RGBColor(255, 193, 7),      # Golden yellow
            "text": RGBColor(33, 37, 41),         # Dark gray
            "success": RGBColor(40, 167, 69)      # Success green
        }
        super().__init__("Marketing", colors)
    
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        
        return f"""
Create a MARKETING PRESENTATION for {company_profile.name}. Focus on benefits and transformation success.

COMPANY CONTEXT:
- Industry: {company_profile.industry}
- Business Model: {company_profile.business_model}

USE CASES FOR BENEFITS:
{self._format_use_cases_for_marketing(use_cases)}

Create a persuasive presentation focused on transformation benefits and success potential.

OUTPUT JSON:
{{
    "presentation_title": "Transform Your Business - {company_profile.name}",
    "slides": [
        {{
            "type": "title",
            "title": "Transform Your Business Potential",
            "subtitle": "Digital Success with {company_profile.name}"
        }},
        {{
            "type": "problem",
            "title": "The Challenge Every Business Faces",
            "pain_points": [
                "Current operational inefficiencies limiting growth",
                "Competitive market pressures and rising costs",
                "Technology gaps preventing optimal performance"
            ]
        }},
        {{
            "type": "solution",
            "title": "The Transformation Solution",
            "benefits": [
                "Streamlined operations for maximum efficiency",
                "Enhanced competitive position through innovation",
                "Technology-driven growth and market advantage"
            ]
        }},
        {{
            "type": "results",
            "title": "Real Results You Can Achieve",
            "outcomes": [
                "Operational efficiency: 25-40% improvement potential",
                "Cost optimization: 15-30% reduction opportunities",
                "Revenue growth: Enhanced market positioning"
            ]
        }},
        {{
            "type": "success_stories",
            "title": "Success Stories and Proven Results",
            "stories": [
                "Manufacturing efficiency gains: 35% improvement",
                "Service delivery optimization: 50% faster processing",
                "Customer satisfaction enhancement: 40% increase"
            ]
        }},
        {{
            "type": "call_to_action",
            "title": "Start Your Transformation Journey",
            "actions": [
                "Schedule your free transformation assessment",
                "Join pilot program for early adopters",
                "Explore strategic partnership opportunities"
            ]
        }}
    ]
}}

Focus on emotional connection and clear benefits.
"""
    
    def _format_use_cases_for_marketing(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for marketing benefits"""
        if not use_cases:
            return "Standard transformation benefits"
        
        formatted = []
        for uc in use_cases[:3]:
            formatted.append(f"• {uc.title}: {uc.business_value[:80]}...")
        
        return "\n".join(formatted)
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create marketing slides"""
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_marketing_title(ppt, slide_data)
            elif slide_type == 'problem':
                self._create_problem_slide(ppt, slide_data)
            elif slide_type == 'solution':
                self._create_solution_slide(ppt, slide_data)
            elif slide_type == 'results':
                self._create_results_slide(ppt, slide_data)
            elif slide_type == 'success_stories':
                self._create_success_slide(ppt, slide_data)
            elif slide_type == 'call_to_action':
                self._create_cta_slide(ppt, slide_data)
    
    def _create_marketing_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Transform Your Business')
        subtitle.text = data.get('subtitle', 'Digital Success')
        
        # Marketing styling - bold and engaging
        title.text_frame.paragraphs[0].font.size = Pt(52)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(30)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
        subtitle.text_frame.paragraphs[0].font.bold = True
    
    def _create_problem_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'The Challenge')
        
        pain_points = data.get('pain_points', [])
        if pain_points:
            content.text = pain_points[0]
            for point in pain_points[1:]:
                p = content.text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        self._style_marketing_slide(title, content, self.colors["accent"])
    
    def _create_solution_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'The Solution')
        
        benefits = data.get('benefits', [])
        if benefits:
            content.text = benefits[0]
            for benefit in benefits[1:]:
                p = content.text_frame.add_paragraph()
                p.text = benefit
                p.level = 0
        
        self._style_marketing_slide(title, content, self.colors["success"])
    
    def _create_results_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Real Results')
        
        outcomes = data.get('outcomes', [])
        if outcomes:
            content.text = outcomes[0]
            for outcome in outcomes[1:]:
                p = content.text_frame.add_paragraph()
                p.text = outcome
                p.level = 0
        
        self._style_marketing_slide(title, content, self.colors["success"])
    
    def _create_success_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Success Stories')
        
        stories = data.get('stories', [])
        if stories:
            content.text = stories[0]
            for story in stories[1:]:
                p = content.text_frame.add_paragraph()
                p.text = story
                p.level = 0
        
        self._style_marketing_slide(title, content, self.colors["text"])
    
    def _create_cta_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Take Action')
        
        actions = data.get('actions', [])
        if actions:
            content.text = actions[0]
            for action in actions[1:]:
                p = content.text_frame.add_paragraph()
                p.text = action
                p.level = 0
        
        self._style_marketing_slide(title, content, self.colors["accent"])
    
    def _style_marketing_slide(self, title, content, text_color) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = text_color
            paragraph.space_before = Pt(10)
            paragraph.font.bold = True

class UseCaseTemplate(PPTTemplate):
    """Use Case Template - Detailed scenarios"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(99, 102, 241),    # Indigo
            "secondary": RGBColor(139, 69, 19),   # Brown
            "accent": RGBColor(245, 158, 11),     # Amber
            "text": RGBColor(55, 65, 81),         # Gray
            "background": RGBColor(249, 250, 251) # Light
        }
        super().__init__("Use Case", colors)
    
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        
        return f"""
Create a detailed USE CASE PRESENTATION for {company_profile.name}.

COMPANY CONTEXT:
- Industry: {company_profile.industry}
- Business Model: {company_profile.business_model}

DETAILED USE CASES:
{self._format_detailed_use_cases(use_cases)}

Create implementation-focused presentation with specific scenarios.

OUTPUT JSON:
{{
    "presentation_title": "{company_profile.name} Use Case Implementation Guide",
    "slides": [
        {{
            "type": "title",
            "title": "Use Case Implementation Strategy",
            "subtitle": "{company_profile.name} Transformation Scenarios"
        }},
        {{
            "type": "overview",
            "title": "Use Case Portfolio Overview",
            "use_cases": [
                "Use Case 1: Process automation and efficiency optimization",
                "Use Case 2: Data analytics and business intelligence enhancement", 
                "Use Case 3: Customer experience and operational excellence"
            ]
        }},
        {{
            "type": "detailed_use_case_1",
            "title": "Use Case 1: Process Optimization",
            "problem": "Current manual processes causing delays and inefficiencies",
            "solution": "Automated workflow management and intelligent task routing",
            "benefits": ["40% reduction in processing time", "60% decrease in manual errors", "$200K annual cost savings"],
            "timeline": "6-8 months implementation"
        }},
        {{
            "type": "detailed_use_case_2",
            "title": "Use Case 2: Data Analytics Enhancement",
            "problem": "Limited visibility into business performance and trends",
            "solution": "Real-time analytics dashboard and predictive insights platform",
            "benefits": ["50% faster decision-making", "25% improvement in forecast accuracy", "$150K efficiency gains"],
            "timeline": "4-6 months implementation"
        }},
        {{
            "type": "detailed_use_case_3",
            "title": "Use Case 3: Customer Experience Optimization",
            "problem": "Inconsistent customer service and lengthy response times",
            "solution": "Integrated customer experience platform with AI-powered support",
            "benefits": ["35% faster response times", "60% improvement in satisfaction", "$300K revenue impact"],
            "timeline": "8-10 months implementation"
        }},
        {{
            "type": "implementation",
            "title": "Implementation Methodology",
            "phases": [
                "Phase 1: Discovery and requirements gathering (months 1-2)",
                "Phase 2: Solution design and pilot development (months 3-4)",
                "Phase 3: Full deployment and optimization (months 5-8)"
            ]
        }},
        {{
            "type": "risk_management",
            "title": "Risk Management and Success Factors",
            "risks": ["Technical integration challenges", "User adoption resistance", "Data migration complexity"],
            "mitigations": ["Phased implementation approach", "Comprehensive training program", "Dedicated project management"]
        }}
    ]
}}

Focus on practical implementation details and realistic timelines.
"""
    
    def _format_detailed_use_cases(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for detailed analysis"""
        if not use_cases:
            return "Standard implementation scenarios"
        
        formatted = []
        for i, uc in enumerate(use_cases[:3], 1):
            formatted.append(f"""
Use Case {i}: {uc.title}
- Current State: {uc.current_state[:100]}...
- Solution: {uc.proposed_solution[:100]}...
- Timeline: {uc.timeline_months} months
""")
        
        return "\n".join(formatted)
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create use case slides"""
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_usecase_title(ppt, slide_data)
            elif slide_type == 'overview':
                self._create_overview_slide(ppt, slide_data)
            elif slide_type.startswith('detailed_use_case'):
                self._create_detailed_slide(ppt, slide_data)
            elif slide_type == 'implementation':
                self._create_implementation_slide(ppt, slide_data)
            elif slide_type == 'risk_management':
                self._create_risk_slide(ppt, slide_data)
    
    def _create_usecase_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Use Case Strategy')
        subtitle.text = data.get('subtitle', 'Implementation Guide')
        
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(26)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_overview_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Use Case Portfolio')
        
        use_cases = data.get('use_cases', [])
        if use_cases:
            content.text = use_cases[0]
            for uc in use_cases[1:]:
                p = content.text_frame.add_paragraph()
                p.text = uc
                p.level = 0
        
        self._style_usecase_slide(title, content)
    
    def _create_detailed_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Detailed Use Case')
        
        problem = data.get('problem', '')
        solution = data.get('solution', '')
        benefits = data.get('benefits', [])
        timeline = data.get('timeline', '')
        
        content_text = f"Problem: {problem}\n\nSolution: {solution}\n\nBenefits:\n"
        for benefit in benefits:
            content_text += f"• {benefit}\n"
        content_text += f"\nTimeline: {timeline}"
        
        content.text = content_text
        self._style_usecase_slide(title, content)
    
    def _create_implementation_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Implementation')
        
        phases = data.get('phases', [])
        if phases:
            content.text = phases[0]
            for phase in phases[1:]:
                p = content.text_frame.add_paragraph()
                p.text = phase
                p.level = 0
        
        self._style_usecase_slide(title, content)
    
    def _create_risk_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Risk Management')
        
        risks = data.get('risks', [])
        mitigations = data.get('mitigations', [])
        
        content_text = "Risk Factors:\n"
        for risk in risks:
            content_text += f"• {risk}\n"
        
        content_text += "\nMitigation Strategies:\n"
        for mitigation in mitigations:
            content_text += f"• {mitigation}\n"
        
        content.text = content_text
        self._style_usecase_slide(title, content)
    
    def _style_usecase_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(6)

class TechnicalTemplate(PPTTemplate):
    """Technical Template - Architecture focused"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(30, 41, 59),      # Slate
            "secondary": RGBColor(71, 85, 105),   # Gray
            "accent": RGBColor(14, 165, 233),     # Blue
            "text": RGBColor(51, 65, 85),         # Dark gray
            "code": RGBColor(239, 68, 68)         # Red
        }
        super().__init__("Technical", colors)
    
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        
        return f"""
Create a TECHNICAL PRESENTATION for {company_profile.name}.

TECHNICAL USE CASES:
{self._format_technical_use_cases(use_cases)}

Create architecture-focused technical presentation.

OUTPUT JSON:
{{
    "presentation_title": "{company_profile.name} Technical Architecture",
    "slides": [
        {{
            "type": "title",
            "title": "Technical Architecture Overview",
            "subtitle": "{company_profile.name} System Design and Implementation"
        }},
        {{
            "type": "architecture",
            "title": "System Architecture Components",
            "components": [
                "Frontend layer: User interface and experience design",
                "Backend services: Business logic and API management",
                "Data layer: Storage, processing, and analytics platform",
                "Integration layer: External system connectivity"
            ]
        }},
        {{
            "type": "technology_stack",
            "title": "Technology Stack and Platforms",
            "technologies": [
                "Frontend: React/Vue.js with TypeScript and responsive design",
                "Backend: Node.js/Python microservices with REST/GraphQL APIs",
                "Database: PostgreSQL/MongoDB with Redis caching layer",
                "Cloud: AWS/Azure with container orchestration (Kubernetes)"
            ]
        }},
        {{
            "type": "integration",
            "title": "Integration Architecture",
            "patterns": [
                "API-first design approach with RESTful services",
                "Event-driven architecture for real-time processing",
                "Microservices pattern for scalable deployment"
            ]
        }},
        {{
            "type": "performance",
            "title": "Performance and Security",
            "requirements": [
                "Performance: <200ms API response, 99.9% uptime SLA",
                "Security: End-to-end encryption, multi-factor authentication",
                "Scalability: Auto-scaling, load balancing, CDN optimization"
            ]
        }}
    ]
}}

Focus on technical specifications and implementation details.
"""
    
    def _format_technical_use_cases(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for technical focus"""
        if not use_cases:
            return "Standard technical implementation"
        
        formatted = []
        for uc in use_cases[:3]:
            formatted.append(f"• {uc.title}: Technical implementation with {uc.technology_stack if hasattr(uc, 'technology_stack') else 'modern architecture'}")
        
        return "\n".join(formatted)
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create technical slides"""
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_tech_title(ppt, slide_data)
            elif slide_type == 'architecture':
                self._create_architecture_slide(ppt, slide_data)
            elif slide_type == 'technology_stack':
                self._create_tech_stack_slide(ppt, slide_data)
            elif slide_type == 'integration':
                self._create_integration_slide(ppt, slide_data)
            elif slide_type == 'performance':
                self._create_performance_slide(ppt, slide_data)
    
    def _create_tech_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Technical Architecture')
        subtitle.text = data.get('subtitle', 'System Design')
        
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_architecture_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'System Architecture')
        
        components = data.get('components', [])
        if components:
            content.text = components[0]
            for comp in components[1:]:
                p = content.text_frame.add_paragraph()
                p.text = comp
                p.level = 0
        
        self._style_tech_slide(title, content)
    
    def _create_tech_stack_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Technology Stack')
        
        technologies = data.get('technologies', [])
        if technologies:
            content.text = technologies[0]
            for tech in technologies[1:]:
                p = content.text_frame.add_paragraph()
                p.text = tech
                p.level = 0
        
        self._style_tech_slide(title, content)
    
    def _create_integration_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Integration Architecture')
        
        patterns = data.get('patterns', [])
        if patterns:
            content.text = patterns[0]
            for pattern in patterns[1:]:
                p = content.text_frame.add_paragraph()
                p.text = pattern
                p.level = 0
        
        self._style_tech_slide(title, content)
    
    def _create_performance_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Performance & Security')
        
        requirements = data.get('requirements', [])
        if requirements:
            content.text = requirements[0]
            for req in requirements[1:]:
                p = content.text_frame.add_paragraph()
                p.text = req
                p.level = 0
        
        self._style_tech_slide(title, content)
    
    def _style_tech_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(34)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(4)

class StrategyTemplate(PPTTemplate):
    """Strategy Template - Strategic planning focused"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(79, 70, 229),     # Indigo
            "secondary": RGBColor(107, 114, 128), # Gray
            "accent": RGBColor(16, 185, 129),     # Emerald
            "text": RGBColor(17, 24, 39),         # Dark
            "highlight": RGBColor(245, 101, 101) # Red
        }
        super().__init__("Strategy", colors)
    
    def get_analysis_prompt(self, company_profile: CompanyProfile, research_data: Dict[str, Any], 
                           use_cases: List[UseCaseStructured]) -> str:
        
        return f"""
Create a STRATEGIC PLANNING presentation for {company_profile.name}.

STRATEGIC USE CASES:
{self._format_strategic_use_cases(use_cases)}

Create strategic roadmap presentation.

OUTPUT JSON:
{{
    "presentation_title": "{company_profile.name} Strategic Transformation Plan",
    "slides": [
        {{
            "type": "title",
            "title": "Strategic Transformation Roadmap",
            "subtitle": "{company_profile.name} 3-Year Strategic Plan"
        }},
        {{
            "type": "current_state",
            "title": "Current State Assessment",
            "strengths": [
                "Market leadership in core business segments",
                "Strong customer relationships and brand recognition",
                "Experienced leadership team and skilled workforce"
            ],
            "challenges": [
                "Technology infrastructure modernization needs",
                "Operational efficiency optimization opportunities",
                "Market expansion and competitive positioning requirements"
            ]
        }},
        {{
            "type": "strategic_vision",
            "title": "Strategic Vision and Objectives",
            "vision": "To become the industry leader in innovation and customer satisfaction through strategic transformation",
            "objectives": [
                "Achieve 25% market share growth over 3 years",
                "Implement digital transformation across all operations",
                "Establish new revenue streams and strategic partnerships"
            ]
        }},
        {{
            "type": "strategic_initiatives",
            "title": "Strategic Initiatives Portfolio",
            "initiatives": [
                "Growth Initiatives: Market expansion and new product development",
                "Efficiency Initiatives: Process optimization and cost reduction",
                "Innovation Initiatives: Technology advancement and R&D investment"
            ]
        }},
        {{
            "type": "roadmap",
            "title": "3-Year Implementation Roadmap",
            "year_1": "Foundation building: Infrastructure modernization and process optimization",
            "year_2": "Growth acceleration: Market expansion and new product launches",
            "year_3": "Market leadership: Innovation commercialization and competitive advantage"
        }},
        {{
            "type": "success_metrics",
            "title": "Success Metrics and Governance",
            "metrics": [
                "Financial: 25% revenue growth, 15% margin improvement",
                "Operational: 30% efficiency gains, 90% customer satisfaction",
                "Strategic: Market leadership position, innovation pipeline"
            ]
        }}
    ]
}}

Focus on strategic planning and long-term roadmaps.
"""
    
    def _format_strategic_use_cases(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for strategic planning"""
        if not use_cases:
            return "Standard strategic initiatives"
        
        formatted = []
        for uc in use_cases[:3]:
            strategic_impact = getattr(uc, 'strategic_impact', 'strategic business transformation')
            formatted.append(f"• {uc.title}: {strategic_impact}")
        
        return "\n".join(formatted)
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create strategy slides"""
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_strategy_title(ppt, slide_data)
            elif slide_type == 'current_state':
                self._create_current_state_slide(ppt, slide_data)
            elif slide_type == 'strategic_vision':
                self._create_vision_slide(ppt, slide_data)
            elif slide_type == 'strategic_initiatives':
                self._create_initiatives_slide(ppt, slide_data)
            elif slide_type == 'roadmap':
                self._create_roadmap_slide(ppt, slide_data)
            elif slide_type == 'success_metrics':
                self._create_metrics_slide(ppt, slide_data)
    
    def _create_strategy_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Plan')
        subtitle.text = data.get('subtitle', '3-Year Roadmap')
        
        title.text_frame.paragraphs[0].font.size = Pt(46)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_current_state_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Current State')
        
        strengths = data.get('strengths', [])
        challenges = data.get('challenges', [])
        
        content_text = "Organizational Strengths:\n"
        for strength in strengths:
            content_text += f"• {strength}\n"
        
        content_text += "\nStrategic Challenges:\n"
        for challenge in challenges:
            content_text += f"• {challenge}\n"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_vision_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Vision')
        
        vision = data.get('vision', '')
        objectives = data.get('objectives', [])
        
        content_text = f"Vision Statement:\n{vision}\n\nStrategic Objectives:\n"
        for obj in objectives:
            content_text += f"• {obj}\n"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_initiatives_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Initiatives')
        
        initiatives = data.get('initiatives', [])
        if initiatives:
            content.text = initiatives[0]
            for initiative in initiatives[1:]:
                p = content.text_frame.add_paragraph()
                p.text = initiative
                p.level = 0
        
        self._style_strategy_slide(title, content)
    
    def _create_roadmap_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Implementation Roadmap')
        
        year_1 = data.get('year_1', '')
        year_2 = data.get('year_2', '')
        year_3 = data.get('year_3', '')
        
        content_text = f"Year 1: {year_1}\n\nYear 2: {year_2}\n\nYear 3: {year_3}"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_metrics_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Success Metrics')
        
        metrics = data.get('metrics', [])
        if metrics:
            content.text = metrics[0]
            for metric in metrics[1:]:
                p = content.text_frame.add_paragraph()
                p.text = metric
                p.level = 0
        
        self._style_strategy_slide(title, content)
    
    def _style_strategy_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(38)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(7)

# Template Registry
TEMPLATE_REGISTRY = {
    "first_deck": FirstDeckTemplate,
    "marketing": MarketingTemplate, 
    "use_case": UseCaseTemplate,
    "technical": TechnicalTemplate,
    "strategy": StrategyTemplate
}

class MultiTemplatePPTGenerator:
    """Multi-template PowerPoint presentation generator integrated with orchestrator"""
    
    def __init__(self, model_manager: EnhancedModelManager):
        self.model_manager = model_manager
        
        if not PPTX_AVAILABLE:
            logger.warning("PowerPoint generation not available - python-pptx not installed")
    
    def generate_presentation(self, company_profile: CompanyProfile, use_cases: List[UseCaseStructured], 
                            research_data: Dict[str, Any], session_id: str, status_tracker: StatusTracker = None,
                            presentation_style: str = "first_deck") -> Optional[str]:
        """Generate PowerPoint presentation using specified template"""
        
        if not PPTX_AVAILABLE:
            logger.error("PowerPoint generation not available")
            return None
        
        # Validate template
        if presentation_style not in TEMPLATE_REGISTRY:
            logger.warning(f"Invalid template {presentation_style}, using first_deck")
            presentation_style = "first_deck"
        
        logger.info(f"Generating {presentation_style} presentation for {company_profile.name}")
        
        try:
            # Update status
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_GENERATING,
                    {'template': presentation_style, 'company': company_profile.name}
                )
            
            # Get template instance
            template = TEMPLATE_REGISTRY[presentation_style]()
            
            # Generate content structure using Bedrock
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_ANALYZING,
                    {'template': presentation_style, 'phase': 'content_analysis'}
                )
            
            structure = self._analyze_with_bedrock(template, company_profile, research_data, use_cases)
            
            # Create PowerPoint presentation
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_CREATING,
                    {'template': presentation_style, 'slides_planned': len(structure.get('slides', []))}
                )
            
            ppt = Presentation()
            template.create_slides(ppt, structure)
            
            # Style and save presentation
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_SAVING,
                    {'template': presentation_style, 'slides_created': len(ppt.slides)}
                )
            
            ppt_filename = self._save_presentation(ppt, company_profile.name, presentation_style, session_id)
            
            # Upload to S3
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_UPLOADING,
                    {'template': presentation_style, 'filename': ppt_filename}
                )
            
            s3_url = self._upload_to_s3(ppt_filename, session_id, company_profile.name, presentation_style)
            
            # Final status update
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PRESENTATION_COMPLETED,
                    {
                        'template': presentation_style,
                        'filename': ppt_filename,
                        's3_url': s3_url,
                        'slides_count': len(ppt.slides)
                    }
                )
            
            logger.info(f"Generated {presentation_style} presentation: {ppt_filename}")
            return s3_url
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.ERROR,
                    {'error': str(e), 'phase': 'presentation_generation'}
                )
            return None
    
    def _analyze_with_bedrock(self, template: PPTTemplate, company_profile: CompanyProfile, 
                             research_data: Dict[str, Any], use_cases: List[UseCaseStructured]) -> Dict[str, Any]:
        """Use Bedrock to analyze content and create presentation structure"""
        
        try:
            # Get template-specific analysis prompt
            prompt = template.get_analysis_prompt(company_profile, research_data, use_cases)
            
            # Use creative model for presentation generation
            bedrock_agent = Agent(model=self.model_manager.creative_model)
            response = bedrock_agent(prompt)
            
            # Parse JSON response
            import re
            json_match = re.search(r'\{.*\}', str(response), re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
            else:
                return self._create_fallback_structure(template.name, company_profile.name)
                
        except Exception as e:
            logger.error(f"Bedrock analysis failed: {e}")
            return self._create_fallback_structure(template.name, company_profile.name)
    
    def _create_fallback_structure(self, template_name: str, company_name: str) -> Dict[str, Any]:
        """Create fallback structure if Bedrock fails"""
        return {
            "presentation_title": f"{company_name} {template_name} Presentation",
            "slides": [
                {
                    "type": "title",
                    "title": f"{company_name} {template_name}",
                    "subtitle": "Business Transformation Overview"
                }
            ]
        }
    
    def _save_presentation(self, ppt: Presentation, company_name: str, template_style: str, session_id: str) -> str:
        """Save presentation to temporary file"""
        
        # Create filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        company_clean = company_name.replace(' ', '_').lower()
        filename = f"{company_clean}_{template_style}_presentation_{timestamp}.pptx"
        
        # Save to temp directory
        tmp_dir = os.path.join(LAMBDA_TMP_DIR, session_id)
        os.makedirs(tmp_dir, exist_ok=True)
        
        filepath = os.path.join(tmp_dir, filename)
        ppt.save(filepath)
        
        return filepath
    
    def _upload_to_s3(self, filepath: str, session_id: str, company_name: str, template_style: str) -> Optional[str]:
        """Upload presentation to S3"""
        
        try:
            filename = os.path.basename(filepath)
            s3_key = f"presentations/{session_id}/{filename}"
            
            # Upload to S3
            s3_client.upload_file(filepath, S3_BUCKET, s3_key)
            
            # Generate URL
            s3_url = f"https://{S3_BUCKET}.s3.amazonaws.com/{s3_key}"
            
            # Cleanup local file
            try:
                os.unlink(filepath)
            except Exception as e:
                logger.warning(f"Failed to cleanup presentation file: {e}")
            
            logger.info(f"Uploaded presentation to S3: {s3_url}")
            return s3_url
            
        except Exception as e:
            logger.error(f"Failed to upload presentation to S3: {e}")
            return None

# Export the main class for use in orchestrator
__all__ = ['MultiTemplatePPTGenerator', 'TEMPLATE_REGISTRY']