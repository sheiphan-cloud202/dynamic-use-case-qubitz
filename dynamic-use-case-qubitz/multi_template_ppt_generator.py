#!/usr/bin/env python3
"""
Multi-Template PowerPoint Generation System
Creates 5 different types of presentations with unique content structures and themes
"""
import os
import boto3
import json
from datetime import datetime
from typing import Dict, Any, List, Optional
from abc import ABC, abstractmethod

# PDF processing
try:
    import PyPDF2
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

class PPTTemplate(ABC):
    """Abstract base class for PowerPoint templates"""
    
    def __init__(self, name: str, colors: Dict[str, RGBColor]):
        self.name = name
        self.colors = colors
    
    @abstractmethod
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        """Generate Bedrock prompt specific to this template type"""
        pass
    
    @abstractmethod
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        """Create slides specific to this template type"""
        pass

class FirstDeckTemplate(PPTTemplate):
    """First Deck Call Template - High-level executive overview"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(20, 33, 61),     # Deep navy
            "secondary": RGBColor(52, 73, 94),   # Slate blue
            "accent": RGBColor(230, 126, 34),    # Orange
            "text": RGBColor(44, 62, 80),        # Dark blue-gray
            "light": RGBColor(236, 240, 241)     # Light gray
        }
        super().__init__("First Deck", colors)
    
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        return f"""
You are creating a FIRST DECK CALL presentation for {company_name}. This is the initial executive presentation to introduce strategic opportunities.

CONTENT TO ANALYZE:
{content[:10000]}

FIRST DECK STRUCTURE REQUIREMENTS:
- High-level executive summary (no technical details)
- Clear business value proposition
- 3 major strategic opportunities (not detailed solutions)
- Financial impact potential (conservative estimates)
- Next steps for deeper exploration

OUTPUT JSON:
{{
    "presentation_info": {{
        "title": "Strategic Opportunities for {company_name}",
        "type": "first_deck",
        "executive_focus": true
    }},
    "slides": [
        {{
            "type": "title",
            "title": "Strategic Partnership Opportunity",
            "subtitle": "{company_name} Business Transformation Overview"
        }},
        {{
            "type": "company_snapshot",
            "title": "Company Overview",
            "key_metrics": ["Revenue/size", "Industry position", "Key challenges"],
            "current_state": "Brief assessment of current situation"
        }},
        {{
            "type": "opportunity_overview",
            "title": "Strategic Opportunities Identified",
            "opportunities": [
                "Primary opportunity area with high-level benefit",
                "Secondary opportunity with business impact",
                "Third opportunity for competitive advantage"
            ]
        }},
        {{
            "type": "value_proposition",
            "title": "Potential Business Value",
            "value_areas": [
                "Cost optimization potential: X-Y%",
                "Efficiency improvement: A-B%", 
                "Revenue growth enablement: C-D%"
            ],
            "timeline": "Initial results in 6-12 months"
        }},
        {{
            "type": "next_steps",
            "title": "Proposed Next Steps",
            "steps": [
                "Detailed assessment and discovery phase",
                "Proof of concept development",
                "Implementation roadmap creation"
            ]
        }}
    ]
}}

Focus on high-level strategic impact, not technical implementation details.
"""
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_title_slide(ppt, slide_data)
            elif slide_type == 'company_snapshot':
                self._create_snapshot_slide(ppt, slide_data)
            elif slide_type == 'opportunity_overview':
                self._create_opportunity_slide(ppt, slide_data)
            elif slide_type == 'value_proposition':
                self._create_value_slide(ppt, slide_data)
            elif slide_type == 'next_steps':
                self._create_next_steps_slide(ppt, slide_data)
    
    def _create_title_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Partnership')
        subtitle.text = data.get('subtitle', 'Executive Overview')
        
        # Executive styling - large, bold, authoritative
        title.text_frame.paragraphs[0].font.size = Pt(48)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_snapshot_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Company Overview')
        
        # High-level metrics focus
        metrics = data.get('key_metrics', [])
        current_state = data.get('current_state', '')
        
        content_text = f"Current State: {current_state}\n\nKey Business Metrics:"
        for metric in metrics:
            content_text += f"\n• {metric}"
        
        content.text = content_text
        
        self._style_content_slide(title, content)
    
    def _create_opportunity_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Opportunities')
        
        opportunities = data.get('opportunities', [])
        if opportunities:
            content.text = opportunities[0]
            for opp in opportunities[1:]:
                p = content.text_frame.add_paragraph()
                p.text = opp
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _create_value_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Business Value')
        
        value_areas = data.get('value_areas', [])
        timeline = data.get('timeline', '')
        
        content_text = "Potential Value Creation:"
        for value in value_areas:
            content_text += f"\n• {value}"
        
        if timeline:
            content_text += f"\n\nTimeline: {timeline}"
        
        content.text = content_text
        self._style_content_slide(title, content)
    
    def _create_next_steps_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Next Steps')
        
        steps = data.get('steps', [])
        if steps:
            content.text = steps[0]
            for step in steps[1:]:
                p = content.text_frame.add_paragraph()
                p.text = step
                p.level = 0
        
        self._style_content_slide(title, content)
    
    def _style_content_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(8)

class MarketingTemplate(PPTTemplate):
    """Marketing Template - Persuasive, benefit-focused, visually engaging"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(225, 45, 139),    # Vibrant pink
            "secondary": RGBColor(74, 144, 226),  # Bright blue
            "accent": RGBColor(255, 193, 7),      # Golden yellow
            "text": RGBColor(33, 37, 41),         # Dark gray
            "success": RGBColor(40, 167, 69)      # Success green
        }
        super().__init__("Marketing", colors)
    
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        return f"""
Create a MARKETING PRESENTATION for {company_name}. Focus on benefits, transformation success, and compelling value propositions.

CONTENT TO ANALYZE:
{content[:10000]}

MARKETING PRESENTATION STRUCTURE:
- Compelling problem statement with emotional impact
- Clear solution narrative with transformation story
- Multiple success scenarios and benefits
- Social proof and credibility elements
- Strong call-to-action

OUTPUT JSON:
{{
    "presentation_info": {{
        "title": "Transform Your Business with {company_name}",
        "type": "marketing",
        "persuasive_focus": true
    }},
    "slides": [
        {{
            "type": "title",
            "title": "Unlock Your Business Potential",
            "subtitle": "Digital Transformation Success with {company_name}"
        }},
        {{
            "type": "problem_agitation",
            "title": "The Challenge Every Business Faces",
            "pain_points": [
                "Specific business pain point with emotional impact",
                "Cost of inaction or competitive disadvantage",
                "Frustrations with current processes"
            ],
            "urgency": "Why immediate action is needed"
        }},
        {{
            "type": "solution_story",
            "title": "The Transformation Solution",
            "story_elements": [
                "How transformation addresses core challenges",
                "Unique approach and differentiators",
                "Technology enablement and innovation"
            ]
        }},
        {{
            "type": "benefits_showcase",
            "title": "Real Results You Can Achieve",
            "benefit_categories": [
                "Operational Excellence: Specific improvements",
                "Financial Impact: Cost savings and revenue growth",
                "Competitive Advantage: Market positioning benefits"
            ]
        }},
        {{
            "type": "success_scenarios",
            "title": "Success Stories and Use Cases",
            "scenarios": [
                "Scenario 1: Department/process transformation",
                "Scenario 2: Technology innovation success",
                "Scenario 3: Market expansion achievement"
            ]
        }},
        {{
            "type": "call_to_action",
            "title": "Start Your Transformation Journey",
            "action_items": [
                "Schedule transformation assessment",
                "Pilot program opportunity",
                "Partnership discussion"
            ]
        }}
    ]
}}

Focus on emotional connection, clear benefits, and persuasive messaging.
"""
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_marketing_title(ppt, slide_data)
            elif slide_type == 'problem_agitation':
                self._create_problem_slide(ppt, slide_data)
            elif slide_type == 'solution_story':
                self._create_solution_slide(ppt, slide_data)
            elif slide_type == 'benefits_showcase':
                self._create_benefits_slide(ppt, slide_data)
            elif slide_type == 'success_scenarios':
                self._create_scenarios_slide(ppt, slide_data)
            elif slide_type == 'call_to_action':
                self._create_cta_slide(ppt, slide_data)
    
    def _create_marketing_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Transform Your Business')
        subtitle.text = data.get('subtitle', 'Digital Success Story')
        
        # Marketing styling - bold, attention-grabbing
        title.text_frame.paragraphs[0].font.size = Pt(52)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(30)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
        subtitle.text_frame.paragraphs[0].font.bold = True
    
    def _create_problem_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'The Challenge')
        
        pain_points = data.get('pain_points', [])
        urgency = data.get('urgency', '')
        
        content_text = ""
        for point in pain_points:
            content_text += f"• {point}\n"
        
        if urgency:
            content_text += f"\n⚠️ {urgency}"
        
        content.text = content_text
        self._style_marketing_slide(title, content, accent_color=True)
    
    def _create_solution_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'The Solution')
        
        story_elements = data.get('story_elements', [])
        if story_elements:
            content.text = story_elements[0]
            for element in story_elements[1:]:
                p = content.text_frame.add_paragraph()
                p.text = element
                p.level = 0
        
        self._style_marketing_slide(title, content)
    
    def _create_benefits_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Real Results')
        
        benefit_categories = data.get('benefit_categories', [])
        if benefit_categories:
            content.text = benefit_categories[0]
            for benefit in benefit_categories[1:]:
                p = content.text_frame.add_paragraph()
                p.text = benefit
                p.level = 0
        
        self._style_marketing_slide(title, content, success_color=True)
    
    def _create_scenarios_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Success Stories')
        
        scenarios = data.get('scenarios', [])
        if scenarios:
            content.text = scenarios[0]
            for scenario in scenarios[1:]:
                p = content.text_frame.add_paragraph()
                p.text = scenario
                p.level = 0
        
        self._style_marketing_slide(title, content)
    
    def _create_cta_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Take Action Now')
        
        action_items = data.get('action_items', [])
        if action_items:
            content.text = action_items[0]
            for item in action_items[1:]:
                p = content.text_frame.add_paragraph()
                p.text = item
                p.level = 0
        
        self._style_marketing_slide(title, content, accent_color=True)
    
    def _style_marketing_slide(self, title, content, accent_color=False, success_color=False) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        text_color = self.colors["accent"] if accent_color else (self.colors["success"] if success_color else self.colors["text"])
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = text_color
            paragraph.space_before = Pt(10)
            if accent_color or success_color:
                paragraph.font.bold = True

class UseCaseTemplate(PPTTemplate):
    """Use Case Template - Detailed problem-solution-benefit structure"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(99, 102, 241),    # Indigo
            "secondary": RGBColor(139, 69, 19),   # Saddle brown
            "accent": RGBColor(245, 158, 11),     # Amber
            "text": RGBColor(55, 65, 81),         # Gray-700
            "background": RGBColor(249, 250, 251) # Light background
        }
        super().__init__("Use Case", colors)
    
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        return f"""
Create a detailed USE CASE PRESENTATION for {company_name}. Focus on specific scenarios, implementation details, and measurable outcomes.

CONTENT TO ANALYZE:
{content[:10000]}

USE CASE STRUCTURE REQUIREMENTS:
- 3-4 specific use cases with detailed problem-solution-benefit
- Implementation methodology and timeline
- Technical considerations and requirements
- Success metrics and measurement approach
- Risk mitigation and change management

OUTPUT JSON:
{{
    "presentation_info": {{
        "title": "{company_name} Use Case Implementation Guide",
        "type": "use_case",
        "detailed_focus": true
    }},
    "slides": [
        {{
            "type": "title",
            "title": "Use Case Implementation Strategy",
            "subtitle": "{company_name} Transformation Scenarios"
        }},
        {{
            "type": "use_case_overview",
            "title": "Use Case Portfolio Overview",
            "use_case_summary": [
                "Use Case 1: Process automation and efficiency",
                "Use Case 2: Data analytics and insights", 
                "Use Case 3: Customer experience enhancement",
                "Use Case 4: Operational cost optimization"
            ]
        }},
        {{
            "type": "detailed_use_case",
            "use_case_number": 1,
            "title": "Use Case 1: [Specific Process Name]",
            "current_state": "Detailed description of current challenge",
            "proposed_solution": "Specific technology and process solution",
            "implementation_steps": [
                "Phase 1: Assessment and planning (months 1-2)",
                "Phase 2: Pilot implementation (months 3-4)",
                "Phase 3: Full rollout (months 5-8)"
            ],
            "expected_benefits": [
                "Quantified efficiency improvement",
                "Cost reduction estimate",
                "Quality improvement metric"
            ],
            "success_metrics": "How success will be measured"
        }},
        {{
            "type": "implementation_approach",
            "title": "Implementation Methodology",
            "methodology_phases": [
                "Discovery and requirements gathering",
                "Solution design and architecture",
                "Pilot development and testing",
                "Training and change management",
                "Full deployment and optimization"
            ],
            "timeline": "12-18 month implementation cycle"
        }},
        {{
            "type": "risk_mitigation",
            "title": "Risk Management and Success Factors",
            "risk_factors": [
                "Technical integration challenges",
                "User adoption and change resistance",
                "Data migration and quality issues"
            ],
            "mitigation_strategies": [
                "Phased implementation approach",
                "Comprehensive training program",
                "Dedicated project management office"
            ]
        }}
    ]
}}

Focus on practical implementation details and realistic timelines.
"""
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_usecase_title(ppt, slide_data)
            elif slide_type == 'use_case_overview':
                self._create_overview_slide(ppt, slide_data)
            elif slide_type == 'detailed_use_case':
                self._create_detailed_usecase(ppt, slide_data)
            elif slide_type == 'implementation_approach':
                self._create_implementation_slide(ppt, slide_data)
            elif slide_type == 'risk_mitigation':
                self._create_risk_slide(ppt, slide_data)
    
    def _create_usecase_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Use Case Strategy')
        subtitle.text = data.get('subtitle', 'Implementation Guide')
        
        # Professional, technical styling
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(26)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_overview_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Use Case Overview')
        
        use_cases = data.get('use_case_summary', [])
        if use_cases:
            content.text = use_cases[0]
            for uc in use_cases[1:]:
                p = content.text_frame.add_paragraph()
                p.text = uc
                p.level = 0
        
        self._style_usecase_slide(title, content)
    
    def _create_detailed_usecase(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Detailed Use Case')
        
        # Structure detailed use case content
        current_state = data.get('current_state', '')
        solution = data.get('proposed_solution', '')
        steps = data.get('implementation_steps', [])
        benefits = data.get('expected_benefits', [])
        metrics = data.get('success_metrics', '')
        
        content_text = f"Current State: {current_state}\n\n"
        content_text += f"Solution: {solution}\n\n"
        
        if steps:
            content_text += "Implementation:\n"
            for step in steps:
                content_text += f"• {step}\n"
        
        content.text = content_text
        self._style_usecase_slide(title, content)
    
    def _create_implementation_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Implementation Approach')
        
        phases = data.get('methodology_phases', [])
        timeline = data.get('timeline', '')
        
        content_text = ""
        for i, phase in enumerate(phases, 1):
            content_text += f"{i}. {phase}\n"
        
        if timeline:
            content_text += f"\nTimeline: {timeline}"
        
        content.text = content_text
        self._style_usecase_slide(title, content)
    
    def _create_risk_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Risk Management')
        
        risks = data.get('risk_factors', [])
        mitigations = data.get('mitigation_strategies', [])
        
        content_text = "Risk Factors:\n"
        for risk in risks:
            content_text += f"• {risk}\n"
        
        content_text += "\nMitigation Strategies:\n"
        for mitigation in mitigations:
            content_text += f"• {mitigation}\n"
        
        content.text = content_text
        self._style_usecase_slide(title, content)
    
    def _style_usecase_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(6)

class TechnicalTemplate(PPTTemplate):
    """Technical Template - Architecture, specifications, implementation details"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(30, 41, 59),      # Slate-800
            "secondary": RGBColor(71, 85, 105),   # Slate-600
            "accent": RGBColor(14, 165, 233),     # Sky-500
            "text": RGBColor(51, 65, 85),         # Slate-700
            "code": RGBColor(239, 68, 68)         # Red-500
        }
        super().__init__("Technical", colors)
    
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        return f"""
Create a TECHNICAL PRESENTATION for {company_name}. Focus on architecture, specifications, and implementation details.

CONTENT TO ANALYZE:
{content[:10000]}

TECHNICAL PRESENTATION STRUCTURE:
- System architecture and technical approach
- Technology stack and platform requirements
- Integration specifications and data flow
- Performance requirements and scalability
- Security and compliance considerations

OUTPUT JSON:
{{
    "presentation_info": {{
        "title": "{company_name} Technical Architecture",
        "type": "technical",
        "technical_focus": true
    }},
    "slides": [
        {{
            "type": "title",
            "title": "Technical Architecture Overview",
            "subtitle": "{company_name} System Design and Implementation"
        }},
        {{
            "type": "architecture_overview",
            "title": "System Architecture",
            "architecture_components": [
                "Frontend layer: User interface and experience",
                "Backend services: Business logic and APIs",
                "Data layer: Storage and analytics platform",
                "Integration layer: External system connectivity"
            ],
            "design_principles": [
                "Scalable and modular architecture",
                "Cloud-native deployment approach",
                "Microservices-based design pattern"
            ]
        }},
        {{
            "type": "technology_stack",
            "title": "Technology Stack and Platforms",
            "frontend_tech": ["React/Vue.js", "TypeScript", "Responsive design"],
            "backend_tech": ["Node.js/Python", "REST APIs", "Microservices"],
            "database_tech": ["PostgreSQL/MongoDB", "Redis caching", "Data warehouse"],
            "cloud_platform": ["AWS/Azure", "Container orchestration", "CI/CD pipeline"]
        }},
        {{
            "type": "integration_specs",
            "title": "Integration Architecture",
            "integration_patterns": [
                "API-first design approach",
                "Event-driven architecture",
                "Real-time data synchronization"
            ],
            "data_flow": [
                "Data ingestion and validation",
                "Processing and transformation",
                "Storage and retrieval optimization"
            ]
        }},
        {{
            "type": "performance_security",
            "title": "Performance and Security",
            "performance_requirements": [
                "Response time: <200ms for API calls",
                "Throughput: 10,000+ concurrent users",
                "Availability: 99.9% uptime SLA"
            ],
            "security_measures": [
                "End-to-end encryption",
                "Multi-factor authentication",
                "Regular security audits and penetration testing"
            ]
        }}
    ]
}}

Focus on technical accuracy, implementation details, and system specifications.
"""
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_tech_title(ppt, slide_data)
            elif slide_type == 'architecture_overview':
                self._create_architecture_slide(ppt, slide_data)
            elif slide_type == 'technology_stack':
                self._create_tech_stack_slide(ppt, slide_data)
            elif slide_type == 'integration_specs':
                self._create_integration_slide(ppt, slide_data)
            elif slide_type == 'performance_security':
                self._create_performance_slide(ppt, slide_data)
    
    def _create_tech_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Technical Architecture')
        subtitle.text = data.get('subtitle', 'System Design')
        
        # Technical styling - clean, precise
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_architecture_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'System Architecture')
        
        components = data.get('architecture_components', [])
        principles = data.get('design_principles', [])
        
        content_text = "Architecture Components:\n"
        for comp in components:
            content_text += f"• {comp}\n"
        
        content_text += "\nDesign Principles:\n"
        for principle in principles:
            content_text += f"• {principle}\n"
        
        content.text = content_text
        self._style_tech_slide(title, content)
    
    def _create_tech_stack_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Technology Stack')
        
        frontend = data.get('frontend_tech', [])
        backend = data.get('backend_tech', [])
        database = data.get('database_tech', [])
        cloud = data.get('cloud_platform', [])
        
        content_text = f"Frontend: {', '.join(frontend)}\n\n"
        content_text += f"Backend: {', '.join(backend)}\n\n"
        content_text += f"Database: {', '.join(database)}\n\n"
        content_text += f"Cloud: {', '.join(cloud)}"
        
        content.text = content_text
        self._style_tech_slide(title, content)
    
    def _create_integration_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Integration Architecture')
        
        patterns = data.get('integration_patterns', [])
        flow = data.get('data_flow', [])
        
        content_text = "Integration Patterns:\n"
        for pattern in patterns:
            content_text += f"• {pattern}\n"
        
        content_text += "\nData Flow:\n"
        for step in flow:
            content_text += f"• {step}\n"
        
        content.text = content_text
        self._style_tech_slide(title, content)
    
    def _create_performance_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Performance & Security')
        
        performance = data.get('performance_requirements', [])
        security = data.get('security_measures', [])
        
        content_text = "Performance Requirements:\n"
        for req in performance:
            content_text += f"• {req}\n"
        
        content_text += "\nSecurity Measures:\n"
        for measure in security:
            content_text += f"• {measure}\n"
        
        content.text = content_text
        self._style_tech_slide(title, content)
    
    def _style_tech_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(34)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(4)

class StrategyTemplate(PPTTemplate):
    """Strategy Template - High-level strategic planning and roadmaps"""
    
    def __init__(self):
        colors = {
            "primary": RGBColor(79, 70, 229),     # Indigo-600
            "secondary": RGBColor(107, 114, 128), # Gray-500
            "accent": RGBColor(16, 185, 129),     # Emerald-500
            "text": RGBColor(17, 24, 39),         # Gray-900
            "highlight": RGBColor(245, 101, 101) # Red-400
        }
        super().__init__("Strategy", colors)
    
    def get_analysis_prompt(self, content: str, company_name: str) -> str:
        return f"""
Create a STRATEGIC PLANNING presentation for {company_name}. Focus on high-level strategy, roadmaps, and organizational transformation.

CONTENT TO ANALYZE:
{content[:10000]}

STRATEGY PRESENTATION STRUCTURE:
- Current state analysis and strategic assessment
- Vision and strategic objectives
- Strategic initiatives and priorities
- Implementation roadmap and milestones
- Success metrics and governance

OUTPUT JSON:
{{
    "presentation_info": {{
        "title": "{company_name} Strategic Transformation Plan",
        "type": "strategy",
        "strategic_focus": true
    }},
    "slides": [
        {{
            "type": "title",
            "title": "Strategic Transformation Roadmap",
            "subtitle": "{company_name} 3-Year Strategic Plan"
        }},
        {{
            "type": "current_state",
            "title": "Current State Assessment",
            "strengths": [
                "Key organizational strengths",
                "Competitive advantages", 
                "Market position benefits"
            ],
            "challenges": [
                "Strategic challenges to address",
                "Market pressures and threats",
                "Internal capability gaps"
            ]
        }},
        {{
            "type": "strategic_vision",
            "title": "Strategic Vision and Objectives",
            "vision_statement": "Clear vision for future state",
            "strategic_objectives": [
                "Objective 1: Market leadership and growth",
                "Objective 2: Operational excellence",
                "Objective 3: Innovation and transformation"
            ]
        }},
        {{
            "type": "strategic_initiatives",
            "title": "Strategic Initiatives Portfolio",
            "initiative_categories": [
                "Growth Initiatives: Market expansion and new products",
                "Efficiency Initiatives: Process optimization and cost reduction",
                "Innovation Initiatives: Technology and capability building"
            ],
            "prioritization": "Based on impact, feasibility, and strategic alignment"
        }},
        {{
            "type": "implementation_roadmap",
            "title": "3-Year Implementation Roadmap",
            "year_1": [
                "Foundation building and quick wins",
                "Core capability development",
                "Initial market expansion"
            ],
            "year_2": [
                "Scale successful initiatives",
                "Technology platform enhancement",
                "Operational optimization"
            ],
            "year_3": [
                "Market leadership achievement",
                "Innovation commercialization",
                "Sustainable competitive advantage"
            ]
        }},
        {{
            "type": "governance_metrics",
            "title": "Success Metrics and Governance",
            "success_metrics": [
                "Financial performance indicators",
                "Market share and customer metrics",
                "Operational efficiency measures"
            ],
            "governance_structure": "Strategic oversight and review process"
        }}
    ]
}}

Focus on strategic thinking, long-term planning, and organizational transformation.
"""
    
    def create_slides(self, ppt: Presentation, structure: Dict[str, Any]) -> None:
        for slide_data in structure.get('slides', []):
            slide_type = slide_data.get('type')
            
            if slide_type == 'title':
                self._create_strategy_title(ppt, slide_data)
            elif slide_type == 'current_state':
                self._create_current_state_slide(ppt, slide_data)
            elif slide_type == 'strategic_vision':
                self._create_vision_slide(ppt, slide_data)
            elif slide_type == 'strategic_initiatives':
                self._create_initiatives_slide(ppt, slide_data)
            elif slide_type == 'implementation_roadmap':
                self._create_roadmap_slide(ppt, slide_data)
            elif slide_type == 'governance_metrics':
                self._create_governance_slide(ppt, slide_data)
    
    def _create_strategy_title(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Plan')
        subtitle.text = data.get('subtitle', '3-Year Roadmap')
        
        # Strategic styling - authoritative, forward-looking
        title.text_frame.paragraphs[0].font.size = Pt(46)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors["secondary"]
    
    def _create_current_state_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Current State')
        
        strengths = data.get('strengths', [])
        challenges = data.get('challenges', [])
        
        content_text = "Strengths:\n"
        for strength in strengths:
            content_text += f"• {strength}\n"
        
        content_text += "\nChallenges:\n"
        for challenge in challenges:
            content_text += f"• {challenge}\n"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_vision_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Vision')
        
        vision = data.get('vision_statement', '')
        objectives = data.get('strategic_objectives', [])
        
        content_text = f"Vision: {vision}\n\nStrategic Objectives:\n"
        for obj in objectives:
            content_text += f"• {obj}\n"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_initiatives_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Strategic Initiatives')
        
        categories = data.get('initiative_categories', [])
        prioritization = data.get('prioritization', '')
        
        content_text = ""
        for category in categories:
            content_text += f"• {category}\n"
        
        if prioritization:
            content_text += f"\nPrioritization: {prioritization}"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_roadmap_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Implementation Roadmap')
        
        year_1 = data.get('year_1', [])
        year_2 = data.get('year_2', [])
        year_3 = data.get('year_3', [])
        
        content_text = "Year 1:\n"
        for item in year_1:
            content_text += f"• {item}\n"
        
        content_text += "\nYear 2:\n"
        for item in year_2:
            content_text += f"• {item}\n"
        
        content_text += "\nYear 3:\n"
        for item in year_3:
            content_text += f"• {item}\n"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _create_governance_slide(self, ppt: Presentation, data: Dict[str, Any]) -> None:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Success Metrics')
        
        metrics = data.get('success_metrics', [])
        governance = data.get('governance_structure', '')
        
        content_text = "Success Metrics:\n"
        for metric in metrics:
            content_text += f"• {metric}\n"
        
        if governance:
            content_text += f"\nGovernance: {governance}"
        
        content.text = content_text
        self._style_strategy_slide(title, content)
    
    def _style_strategy_slide(self, title, content) -> None:
        title.text_frame.paragraphs[0].font.size = Pt(38)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = self.colors["text"]
            paragraph.space_before = Pt(7)

# Template Registry
TEMPLATE_REGISTRY = {
    "first_deck": FirstDeckTemplate,
    "marketing": MarketingTemplate,
    "use_case": UseCaseTemplate,
    "technical": TechnicalTemplate,
    "strategy": StrategyTemplate
}

class PDFContentExtractor:
    """Extract text content from PDF files"""
    
    @staticmethod
    def extract_text(pdf_path: str) -> str:
        if not PDF_AVAILABLE:
            raise ImportError("Install PDF libraries: pip install PyMuPDF PyPDF2")
        
        text_content = ""
        
        try:
            doc = fitz.open(pdf_path)
            for page in doc:
                text_content += page.get_text()
            doc.close()
            
            if len(text_content.strip()) > 100:
                return PDFContentExtractor._clean_text(text_content)
        except Exception:
            pass
        
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content += page.extract_text()
                    
            return PDFContentExtractor._clean_text(text_content)
        except Exception as e:
            raise Exception(f"Could not extract text from PDF: {e}")
    
    @staticmethod
    def _clean_text(text: str) -> str:
        import re
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r' +', ' ', text)
        return text.strip()

class BedrockAnalyzer:
    """Analyze content using AWS Bedrock"""
    
    def __init__(self, region='us-east-1'):
        self.bedrock = boto3.client('bedrock-runtime', region_name=region)
        self.model_id = 'us.anthropic.claude-sonnet-4-20250514-v1:0'
    
    def analyze_with_template(self, content: str, template: PPTTemplate, company_name: str) -> Dict[str, Any]:
        """Analyze content using specific template prompt"""
        
        prompt = template.get_analysis_prompt(content, company_name)
        
        try:
            response = self.bedrock.invoke_model(
                modelId=self.model_id,
                body=json.dumps({
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 4000,
                    "messages": [{"role": "user", "content": prompt}]
                })
            )
            
            response_body = json.loads(response['body'].read())
            ai_response = response_body['content'][0]['text']
            
            # Extract JSON from response
            import re
            json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
            else:
                return self._create_fallback_structure(template.name, company_name)
                
        except Exception as e:
            print(f"Bedrock analysis failed: {e}")
            return self._create_fallback_structure(template.name, company_name)
    
    def _create_fallback_structure(self, template_name: str, company_name: str) -> Dict[str, Any]:
        return {
            "presentation_info": {
                "title": f"{company_name} {template_name} Presentation",
                "type": template_name.lower(),
                "fallback": True
            },
            "slides": [
                {
                    "type": "title",
                    "title": f"{company_name} {template_name}",
                    "subtitle": f"Analysis and Recommendations"
                }
            ]
        }

class MultiTemplatePPTGenerator:
    """Main generator that creates presentations using different templates"""
    
    def __init__(self, aws_region: str = 'us-east-1'):
        self.pdf_extractor = PDFContentExtractor()
        self.bedrock_analyzer = BedrockAnalyzer(aws_region)
        
        if not PPT_AVAILABLE:
            raise ImportError("Install PowerPoint library: pip install python-pptx")
    
    def generate_presentation(self, pdf_path: str, template_type: str, company_name: str = None) -> str:
        """
        Generate presentation using specified template
        
        Args:
            pdf_path: Path to PDF file
            template_type: "first_deck", "marketing", "use_case", "technical", "strategy"  
            company_name: Company name for customization
            
        Returns:
            Path to generated PowerPoint file
        """
        
        # Validate template type
        if template_type not in TEMPLATE_REGISTRY:
            raise ValueError(f"Invalid template type. Choose from: {list(TEMPLATE_REGISTRY.keys())}")
        
        print(f"Generating {template_type} presentation...")
        print(f"PDF: {pdf_path}")
        print(f"Company: {company_name or 'Auto-detect'}")
        
        # Step 1: Extract content from PDF
        print(f"\n1. Extracting content from PDF...")
        content = self.pdf_extractor.extract_text(pdf_path)
        print(f"   Extracted {len(content):,} characters")
        
        # Step 2: Get template and analyze content
        print(f"\n2. Analyzing content with {template_type} template...")
        template = TEMPLATE_REGISTRY[template_type]()
        structure = self.bedrock_analyzer.analyze_with_template(content, template, company_name)
        slide_count = len(structure.get('slides', []))
        print(f"   Generated structure for {slide_count} slides")
        
        # Step 3: Create PowerPoint presentation
        print(f"\n3. Creating {template_type} PowerPoint presentation...")
        ppt = Presentation()
        template.create_slides(ppt, structure)
        
        # Save presentation
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        company_clean = (company_name or 'business').replace(' ', '_').lower()
        filename = f"{company_clean}_{template_type}_presentation_{timestamp}.pptx"
        
        ppt.save(filename)
        
        file_size = os.path.getsize(filename)
        print(f"   Created: {filename} ({file_size:,} bytes)")
        
        return filename

def main():
    """User interface for multi-template presentation generation"""
    
    print("Multi-Template PowerPoint Generator")
    print("=" * 40)
    
    # Check dependencies
    if not PDF_AVAILABLE:
        print("Missing PDF libraries. Install with:")
        print("pip install PyMuPDF PyPDF2")
        return
    
    if not PPT_AVAILABLE:
        print("Missing PowerPoint library. Install with:")
        print("pip install python-pptx")
        return
    
    # Get inputs
    pdf_path = input("Enter PDF file path: ").strip().strip('"\'')
    
    if not os.path.exists(pdf_path):
        print(f"File not found: {pdf_path}")
        return
    
    company_name = input("Company name (optional): ").strip()
    
    print("\nPresentation Templates:")
    print("1. First Deck Call (High-level executive overview)")
    print("2. Marketing Presentation (Persuasive, benefit-focused)")
    print("3. Use Case Scenarios (Detailed problem-solution-benefit)")
    print("4. Technical Architecture (Specifications and implementation)")
    print("5. Strategy Planning (Roadmaps and strategic initiatives)")
    
    choice = input("Choose template (1-5): ").strip()
    template_map = {
        "1": "first_deck",
        "2": "marketing", 
        "3": "use_case",
        "4": "technical",
        "5": "strategy"
    }
    
    template_type = template_map.get(choice)
    if not template_type:
        print("Invalid choice. Using first_deck template.")
        template_type = "first_deck"
    
    try:
        # Generate presentation
        generator = MultiTemplatePPTGenerator()
        result = generator.generate_presentation(pdf_path, template_type, company_name)
        
        print(f"\nSUCCESS!")
        print(f"Generated: {result}")
        print(f"Location: {os.path.abspath(result)}")
        print(f"Template: {template_type}")
        print(f"Ready to edit in PowerPoint")
        
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()